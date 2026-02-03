# Padronização de nomes de arquivos PPTX

from pathlib import Path
from pptx import Presentation
import re


# Configurações
PASTA = Path("CAMINHO/DA/PASTA/COM/PPTX")
EXTENSAO = ".pptx"

SIGLAS_VALIDAS = [
    "FUSO",
    "HESO",
    "INSO",
    "BISO",
    "TSSO",
]


def texto_primeiro_slide(caminho):
    # Lê todo o texto do primeiro slide
    prs = Presentation(caminho)
    slide = prs.slides[0]

    textos = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    textos.append(p.text.strip())

    return " ".join(textos)


def limpar_nome_arquivo(texto, limite=120):
    # Remove caracteres inválidos
    texto = re.sub(r'[\\/:*?"<>|]', '', texto)
    texto = re.sub(r'\s+', ' ', texto).strip()
    texto = re.sub(r'\s+\.', '.', texto)

    return texto[:limite]


def extrair_infos(texto):
    # Sigla
    sigla = "XXXX"
    for s in SIGLAS_VALIDAS:
        if re.search(rf'\b{s}\b', texto):
            sigla = s
            break

    # CODE
    tem_code = bool(re.search(r'\bCODE\b', texto, re.IGNORECASE))

    # Ano e número
    match = re.search(r'\b(20\d{2})[ _-]?(\d{2}[A-B]?)\b', texto, re.IGNORECASE)
    if match:
        ano = match.group(1)
        numero = match.group(2)
        fim_padrao = match.end()
    else:
        ano = "0000"
        numero = "00"
        fim_padrao = 0

    # Título
    titulo = texto[fim_padrao:].strip()
    titulo = re.split(r'\b(Líder|Co líder|Gerente|Versão|Sumário)\b', titulo, flags=re.IGNORECASE)[0]
    titulo = re.sub(r'[-_–]+', ' ', titulo)
    titulo = re.sub(r'^\s*CODE\b\s*[-–]?\s*', '', titulo, flags=re.IGNORECASE)
    titulo = limpar_nome_arquivo(titulo)

    if not titulo:
        titulo = "Sem_Titulo"

    if re.match(r'^[AB]\b', titulo):
        numero = f"{numero}{titulo[0]}"
        titulo = titulo[1:].strip()

    return sigla, ano, numero, tem_code, titulo


print("Iniciando padronização...\n")

for arquivo in sorted(PASTA.glob(f"*{EXTENSAO}")):
    try:
        texto = texto_primeiro_slide(arquivo)
        sigla, ano, numero, tem_code, titulo = extrair_infos(texto)

        codigo = " CODE" if tem_code else ""
        novo_nome = f"{sigla} {ano}_{numero}{codigo} - {titulo}{EXTENSAO}"

        arquivo.rename(arquivo.parent / novo_nome)
        print(f"✔ {novo_nome}")

    except Exception as e:
        print(f"✖ Erro em {arquivo.name}: {e}")

print("\nFinalizado.")
